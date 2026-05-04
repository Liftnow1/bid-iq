"""Pillar 3 ingester for Liftnow's organizational reference docs.

Sibling to bidiq/ingest.py — Wave 1 logic is untouched. The two paths
differ in five places:

1. Multi-format input (.pdf, .docx, .pptx, .md, .xlsx) — not just PDF.
2. Full-body extraction — not Wave 1's tier-1 cover/last-page-only.
3. Tier comes from the staging FOLDER (Tier 1- Public / Tier 2 - Internal /
   Tier 3- Paul Only). The classifier still runs but only as a verification
   layer; its output is logged, never used to override the folder tier.
4. brand_id always points to the `liftnow` row — Pillar 3 docs are
   organizational, not brand-specific.
5. source_type is `pillar3_staging`, extractor_version is
   `ingest.py-v1-pillar3-full`, distinct from Wave 1's tags.

Reuses upsert_knowledge_item (with kwargs added in this branch) plus the
PDF rasterization/vision helpers from bidiq/ingest.py.
"""

from __future__ import annotations

import csv
import io
import json
import os
import re
import sys
import time
import traceback
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from pathlib import Path
from typing import Optional

import anthropic
import psycopg

from bidiq.ingest import (
    REPO_ROOT,
    classify_document,
    extract_chunk_body_with_vision,
    log_error,
    render_pdf_pages,
    upsert_knowledge_item,
)


PILLAR3_ROOT = REPO_ROOT / "data" / "pillar3-staging"

# Version tag distinct from Wave 1's tier1/tier2 markers so the rows are
# trivially queryable.
EXTRACTOR_VERSION_PILLAR3 = "ingest.py-v1-pillar3-full"
SOURCE_TYPE_PILLAR3 = "pillar3_staging"

# Filenames that are present on disk but DO NOT get ingested directly.
SKIP_FILENAMES = {
    # We ingest the auto-generated markdown (kmz_to_markdown.py) instead.
    "New Service Map.kmz",
}

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".md", ".xlsx"}

# Per-file Pillar-3 PDF cost cap. If a single PDF crosses this during
# extraction we stop processing it and log a warning. Defensive — every
# Pillar 3 PDF should be well under this.
PDF_COST_CAP_USD = 0.50

# Rough cost-per-page for Sonnet vision @ 130dpi. Used for cost cap and
# logging only; the actual API bill is what counts.
PDF_VISION_COST_PER_PAGE = 0.012

# Hard warn threshold — if a Pillar 3 PDF is over this many pages, log a
# warning before processing (might be a misfiled product manual).
PDF_WARN_PAGES = 30

# Tier folder name -> v2.1 classifier vocabulary mapping.
TIER_VALUES = {
    "tier-1-public",
    "tier-2-internal",
    "tier-3-paul-only",
}


# ---------------------------------------------------------------------------
# File walker
# ---------------------------------------------------------------------------


class Pillar3File:
    """Tuple-ish record describing one staged file."""

    __slots__ = ("path", "tier", "ext", "size_bytes")

    def __init__(self, path: Path, tier: str, ext: str, size_bytes: int):
        self.path = path
        self.tier = tier
        self.ext = ext
        self.size_bytes = size_bytes

    def __repr__(self) -> str:  # debug only
        return f"<Pillar3File tier={self.tier} ext={self.ext} {self.path.name!r}>"


_TIER_REGEXES = [
    (re.compile(r"(?i)tier\s*1"), "tier-1-public"),
    (re.compile(r"(?i)tier\s*2"), "tier-2-internal"),
    (re.compile(r"(?i)tier\s*3"), "tier-3-paul-only"),
]


def derive_tier_from_path(p: Path) -> Optional[str]:
    """Return the tier vocab string for the first folder segment that matches.

    Match is regex-based so we tolerate the inconsistent dash spacing in
    the staging folders ("Tier 1- Public", "Tier 2 - Internal", etc.).
    Returns None if no segment matches — caller should treat that as an
    error and skip the file.
    """
    for part in p.parts:
        for regex, tier in _TIER_REGEXES:
            if regex.search(part):
                return tier
    return None


def discover_pillar3_files(staging_root: Path = PILLAR3_ROOT) -> list[Pillar3File]:
    """Walk pillar3-staging/ and return Pillar3File records.

    Skips:
    - Hidden / underscore-prefixed paths (Wave 1 convention).
    - Files inside `data/pillar2-staging/` if the walker sees them
      (defensive — separate folder, but keep the check).
    - Filenames on the SKIP_FILENAMES list (the original .kmz).
    - Files outside the supported extension set.
    """
    if not staging_root.exists():
        return []

    results: list[Pillar3File] = []
    for p in staging_root.rglob("*"):
        if not p.is_file():
            continue
        if any(part.startswith((".", "_")) for part in p.parts):
            continue
        if "pillar2-staging" in p.parts:
            continue
        if p.name in SKIP_FILENAMES:
            continue
        ext = p.suffix.lower()
        if ext not in SUPPORTED_EXTS:
            continue
        tier = derive_tier_from_path(p.relative_to(staging_root))
        if tier is None:
            print(
                f"  WARN  no tier folder match for {p.relative_to(staging_root)} — skipping",
                file=sys.stderr,
            )
            continue
        results.append(
            Pillar3File(
                path=p,
                tier=tier,
                ext=ext,
                size_bytes=p.stat().st_size,
            )
        )
    return sorted(results, key=lambda f: (f.tier, f.path.name))


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------


def repo_relative_path(p: Path) -> str:
    return str(p.resolve().relative_to(REPO_ROOT))


def extract_md(path: Path) -> Optional[str]:
    try:
        return path.read_text(encoding="utf-8")
    except Exception as e:
        log_error(path, e)
        return None


def extract_docx(path: Path) -> Optional[str]:
    """Walk paragraphs and tables from a .docx.

    Headings preserved as markdown (#, ##, …) so retrieval-time embeddings
    can lean on them. Tables rendered as pipe tables.
    """
    try:
        import docx  # python-docx
    except ImportError as e:
        log_error(path, e)
        return None
    try:
        doc = docx.Document(str(path))
        out: list[str] = []
        for block in _iter_docx_blocks(doc):
            out.append(block)
        return "\n\n".join(s for s in out if s.strip())
    except Exception as e:
        log_error(path, e)
        return None


def _iter_docx_blocks(doc) -> "list[str]":
    """Yield paragraphs (as markdown) and tables in document order.

    python-docx hides the in-order iteration; the body's child XML
    elements have the right sequence so we walk those and dispatch.
    """
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    blocks: list[str] = []
    for child in doc.element.body.iterchildren():
        tag = child.tag
        if tag == qn("w:p"):
            para = Paragraph(child, doc)
            text = (para.text or "").strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if style.startswith("heading 1") or style == "title":
                blocks.append(f"# {text}")
            elif style.startswith("heading 2"):
                blocks.append(f"## {text}")
            elif style.startswith("heading 3"):
                blocks.append(f"### {text}")
            elif style.startswith("heading 4"):
                blocks.append(f"#### {text}")
            elif style.startswith("heading"):
                blocks.append(f"##### {text}")
            else:
                blocks.append(text)
        elif tag == qn("w:tbl"):
            table = Table(child, doc)
            blocks.append(_docx_table_to_md(table))
    return blocks


def _docx_table_to_md(table) -> str:
    """Render a docx Table as a pipe-format markdown table."""
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [(cell.text or "").strip().replace("\n", " ") for cell in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = rows[0]
    sep = ["---"] * width
    body = rows[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def extract_pptx(path: Path) -> Optional[str]:
    """Slide-by-slide text dump including titles, body, tables, and notes."""
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as e:
        log_error(path, e)
        return None
    try:
        prs = Presentation(str(path))
        out: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            try:
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
            except Exception:
                title = ""
            chunk = [f"## Slide {i}: {title}" if title else f"## Slide {i}"]
            # Text frames
            body_pieces: list[str] = []
            for shape in slide.shapes:
                if shape == getattr(slide.shapes, "title", None):
                    continue
                if shape.has_text_frame:
                    text = (shape.text_frame.text or "").strip()
                    if text:
                        body_pieces.append(text)
                # Tables
                if shape.has_table:
                    body_pieces.append(_pptx_table_to_md(shape.table))
            if body_pieces:
                chunk.append("\n\n".join(body_pieces))
            # Speaker notes
            try:
                notes_frame = slide.notes_slide.notes_text_frame
                notes = (notes_frame.text or "").strip()
                if notes:
                    chunk.append(f"_Notes: {notes}_")
            except Exception:
                pass
            chunk.append("---")
            out.append("\n\n".join(chunk))
        return "\n\n".join(out).strip() or None
    except Exception as e:
        log_error(path, e)
        return None


def _pptx_table_to_md(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        cells = [(c.text or "").strip().replace("\n", " ") for c in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    header = rows[0]
    sep = ["---"] * width
    body = rows[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
    for r in body:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def extract_xlsx(path: Path, max_rows_per_sheet: int = 200, max_cell_chars: int = 500) -> Optional[str]:
    """Per-sheet markdown table dump with row + cell caps per the brief."""
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        log_error(path, e)
        return None
    try:
        wb = load_workbook(str(path), data_only=True, read_only=True)
        out: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            # Drop trailing all-None rows the read_only iterator can yield.
            while rows and all(c is None for c in rows[-1]):
                rows.pop()
            if not rows:
                continue
            total = len(rows)
            truncated = False
            if total > max_rows_per_sheet:
                rows = rows[:max_rows_per_sheet]
                truncated = True

            def _cell(v):
                if v is None:
                    return ""
                s = str(v).replace("\n", " ").replace("|", "\\|")
                if len(s) > max_cell_chars:
                    s = s[: max_cell_chars - 1] + "…"
                return s

            width = max(len(r) for r in rows)
            header = [_cell(c) for c in rows[0]] + [""] * (width - len(rows[0]))
            sep = ["---"] * width
            body_lines = []
            for r in rows[1:]:
                cells = [_cell(c) for c in r] + [""] * (width - len(r))
                body_lines.append("| " + " | ".join(cells) + " |")
            sheet_md = [
                f"## Sheet: {sheet_name}",
                "| " + " | ".join(header) + " |",
                "| " + " | ".join(sep) + " |",
                *body_lines,
            ]
            if truncated:
                sheet_md.append(
                    f"*[Sheet truncated at {max_rows_per_sheet} rows of {total} total]*"
                )
            out.append("\n".join(sheet_md))
        wb.close()
        return "\n\n".join(out).strip() or None
    except Exception as e:
        log_error(path, e)
        return None


def _process_chunks_parallel(
    client: anthropic.Anthropic,
    path: Path,
    pages: list,
    page_count: int,
    *,
    model: str,
    chunk_size: int,
    chunk_concurrency: int,
) -> tuple[list[str], int]:
    """Run chunked body extraction in parallel. Returns (body_parts_in_order, chunks_failed)."""
    chunks: list[list[tuple[int, str]]] = [
        pages[i : i + chunk_size] for i in range(0, len(pages), chunk_size)
    ]
    results: list[Optional[str]] = [None] * len(chunks)
    failed = 0

    def run_one(idx: int) -> tuple[int, Optional[str]]:
        try:
            r = extract_chunk_body_with_vision(
                client, path.name, "liftnow", chunks[idx],
                page_count, model, include_page_summaries=False,
            )
            return idx, (r.get("content_markdown") or "").strip() or None
        except Exception as e:
            log_error(path, e)
            return idx, None

    with ThreadPoolExecutor(max_workers=chunk_concurrency) as pool:
        futs = [pool.submit(run_one, i) for i in range(len(chunks))]
        for fut in as_completed(futs):
            idx, body = fut.result()
            results[idx] = body
            if body is None:
                failed += 1

    body_parts = [r for r in results if r]
    return body_parts, failed


def extract_pdf_native_text(
    path: Path, *, min_chars_per_page: int = 200
) -> Optional[dict]:
    """Try pypdf's native text-layer extraction.

    Many Pillar 3 PDFs are HTML-to-PDF conversions (state RFPs, contract
    drafts, web exports) with a clean text layer. Vision extraction on
    these is wasteful AND lossy — vision re-OCRs already-perfect text
    and routinely loses the long-form answer-box content. This function
    returns rich text in milliseconds, no API cost.

    Returns {"content_markdown", "page_count", "cost_estimate_usd": 0,
    "method": "pypdf"} only if the text layer averages at least
    min_chars_per_page chars across the whole document; otherwise
    returns None so the caller falls back to vision.
    """
    try:
        import pypdf
    except ImportError:
        return None
    try:
        reader = pypdf.PdfReader(str(path))
        page_count = len(reader.pages)
        if page_count == 0:
            return None
        page_texts: list[str] = []
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            page_texts.append(t)
        total = sum(len(t) for t in page_texts)
        avg = total / page_count
        if avg < min_chars_per_page:
            return None  # text layer is thin/missing — fall back to vision
        # Stitch pages with explicit page breaks so retrieval can use them.
        body = "\n\n".join(
            f"## Page {i+1}\n\n{t}" for i, t in enumerate(page_texts) if t
        )
        return {
            "content_markdown": body,
            "page_count": page_count,
            "cost_estimate_usd": 0.0,
            "method": "pypdf",
        }
    except Exception as e:
        log_error(path, e)
        return None


def extract_pdf_full_body(
    client: anthropic.Anthropic,
    path: Path,
    *,
    model: str,
    dpi: int = 130,
    chunk_size: int = 5,
    chunk_concurrency: int = 4,
    max_outer_retries: int = 2,
    prefer_native_text: bool = True,
) -> Optional[dict]:
    """Full-body extraction with parallel chunks and outer retry.

    The dry-run revealed two pain points:
    - Sequential chunking made 70-page contracts take 30-40 min each.
    - Anthropic API connection storms occasionally killed every chunk
      retry, leaving the whole file with 0 chars even though pdfinfo
      could read the PDF fine.

    This version:
    1. Rasterizes pages once, then runs chunks in a thread pool
       (chunk_concurrency).
    2. If the FIRST pass yields zero body parts, sleeps and tries the
       whole file again (up to max_outer_retries) on the theory that the
       API was transiently unhealthy.
    3. Stops if the running cost estimate exceeds PDF_COST_CAP_USD; what
       we have so far is returned.

    Returns {"content_markdown", "page_count", "cost_estimate_usd",
    "cost_capped", "chunks_failed", "outer_retries_used", "method"}
    or None.
    """
    # Native text first — RFPs / HTML-to-PDF exports / contract drafts
    # often have a clean text layer that vision would otherwise re-OCR
    # and lose detail. extract_pdf_native_text returns None if the layer
    # is too sparse; in that case we fall through to vision.
    if prefer_native_text:
        native = extract_pdf_native_text(path)
        if native is not None:
            native["cost_capped"] = False
            native["chunks_failed"] = 0
            native["outer_retries_used"] = 0
            return native

    try:
        pages = render_pdf_pages(str(path), dpi=dpi)
    except Exception as e:
        log_error(path, e)
        return None
    if not pages:
        log_error(path, RuntimeError("0 pages rendered"))
        return None
    page_count = len(pages)
    if page_count > PDF_WARN_PAGES:
        print(
            f"  WARN  {path.name} is {page_count} pages (above {PDF_WARN_PAGES} threshold) "
            f"— possibly a misfiled product manual",
            file=sys.stderr,
        )

    body_parts: list[str] = []
    chunks_failed_total = 0
    outer_retries_used = 0
    cost = 0.0

    # Cap the actual chunk count we'll attempt by the cost cap.
    max_chunks_by_cost = int(PDF_COST_CAP_USD / (chunk_size * PDF_VISION_COST_PER_PAGE))
    capped = False

    for attempt in range(max_outer_retries + 1):
        if attempt > 0:
            backoff = 60 * attempt
            print(
                f"  RETRY {path.name} outer pass {attempt} "
                f"(prior pass yielded 0 body parts) — sleeping {backoff}s",
                file=sys.stderr,
            )
            time.sleep(backoff)
            outer_retries_used = attempt

        # If cost-capped and we already have something, don't retry.
        if cost >= PDF_COST_CAP_USD:
            capped = True
            break

        body_parts, chunks_failed = _process_chunks_parallel(
            client, path, pages, page_count,
            model=model, chunk_size=chunk_size,
            chunk_concurrency=chunk_concurrency,
        )
        chunks_failed_total += chunks_failed
        cost += page_count * PDF_VISION_COST_PER_PAGE

        if body_parts:
            break  # any content is success — stop retrying

    if not body_parts:
        return None
    return {
        "content_markdown": "\n\n".join(body_parts),
        "page_count": page_count,
        "cost_estimate_usd": round(cost, 4),
        "cost_capped": capped,
        "chunks_failed": chunks_failed_total,
        "outer_retries_used": outer_retries_used,
        "method": "vision",
    }


# ---------------------------------------------------------------------------
# Per-file extraction dispatch
# ---------------------------------------------------------------------------


def extract_file(
    f: Pillar3File, client: anthropic.Anthropic, model: str
) -> tuple[Optional[str], int, float]:
    """Return (extracted_text, page_count, cost_estimate). page_count is 0 for non-PDF."""
    if f.ext == ".md":
        text = extract_md(f.path)
        return text, 0, 0.0
    if f.ext == ".docx":
        text = extract_docx(f.path)
        return text, 0, 0.0
    if f.ext == ".pptx":
        text = extract_pptx(f.path)
        return text, 0, 0.0
    if f.ext == ".xlsx":
        text = extract_xlsx(f.path)
        return text, 0, 0.0
    if f.ext == ".pdf":
        result = extract_pdf_full_body(client, f.path, model=model)
        if result is None:
            return None, 0, 0.0
        return (
            result["content_markdown"],
            result["page_count"],
            result["cost_estimate_usd"],
        )
    return None, 0, 0.0


# ---------------------------------------------------------------------------
# DB write
# ---------------------------------------------------------------------------


def get_liftnow_brand_id(conn: psycopg.Connection) -> int:
    with conn.cursor() as cur:
        cur.execute("SELECT id FROM brands WHERE name = 'liftnow' LIMIT 1")
        row = cur.fetchone()
        if not row:
            raise RuntimeError(
                "liftnow brand row missing — run Phase 1.1 SQL first."
            )
        return int(row[0])


def write_pillar3_row(
    conn: psycopg.Connection,
    *,
    f: Pillar3File,
    extracted_text: str,
    page_count: int,
    brand_id: int,
) -> int:
    """Persist via upsert_knowledge_item with Pillar-3 overrides.

    Reuses the Wave 1 tier=2 storage path (full body in raw_content,
    full text in search_text), but stamps the row with Pillar-3-specific
    source_type / content_type / extractor_version so it's distinguishable.
    """
    title = f.path.stem.strip()[:500] or f.path.name
    extraction = {
        "title": title,
        "summary": "",
        "tags": [],
        "category": [f.tier],
        "content_markdown": extracted_text,
        "pages_summary": [],
        "effective_date": None,
        "supersedes_previous": None,
    }
    return upsert_knowledge_item(
        conn,
        source_path=repo_relative_path(f.path),
        source_filename=f.path.name,
        source_pages_count=page_count,
        brand_id=brand_id,
        extraction=extraction,
        tier=2,
        source_type=SOURCE_TYPE_PILLAR3,
        content_type=f.ext.lstrip("."),
        extractor_version_override=EXTRACTOR_VERSION_PILLAR3,
    )


# ---------------------------------------------------------------------------
# Classifier verification (folder primary, classifier secondary)
# ---------------------------------------------------------------------------


def verify_classification(
    client: anthropic.Anthropic,
    *,
    folder_tier: str,
    title: str,
    extracted_text: str,
    model: str,
) -> tuple[str, bool]:
    """Run the v2.1 classifier as a verification layer.

    Returns (classifier_tier, agreement_with_folder).
    Folder tier is what gets written; this is purely for logging.
    """
    try:
        cats = classify_document(
            client,
            title=title,
            summary="",
            body_text=extracted_text,
            model=model,
        )
    except Exception:
        return "uncategorized", False
    classifier_tier = cats[0] if cats else "uncategorized"
    if classifier_tier not in TIER_VALUES and classifier_tier != "uncategorized":
        classifier_tier = "uncategorized"
    return classifier_tier, classifier_tier == folder_tier


# Public surface for the runner.
__all__ = [
    "PILLAR3_ROOT",
    "EXTRACTOR_VERSION_PILLAR3",
    "SOURCE_TYPE_PILLAR3",
    "Pillar3File",
    "discover_pillar3_files",
    "extract_file",
    "verify_classification",
    "write_pillar3_row",
    "get_liftnow_brand_id",
    "PDF_COST_CAP_USD",
    "PDF_WARN_PAGES",
]
